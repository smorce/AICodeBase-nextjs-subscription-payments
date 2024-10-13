import os
import subprocess
import aiofiles
import urllib
import uuid
import mistune
import re
from md2pdf.core import md2pdf
from docx import Document
from htmldocx import HtmlToDocx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from datetime import datetime

# 日本語に翻訳するときに Gemini の出力が約8000文字までなので、6000文字ごとに区切って日本語訳する
def split_text_into_chunks(text: str, max_length: int = 6000) -> list:
    """
    指定されたテキストを最大長指定でチャンクに分割します。

    Parameters:
    text (str): 分割するテキストデータ。
    max_length (int): 各チャンクの最大文字数（デフォルトは6000）。

    Returns:
    list: チャンクに分割されたテキストのリスト。
    """
    lines = text.split('\n')
    chunks = []
    current_chunk = ''
    last_breakpoint = 0

    for i, line in enumerate(lines):
        # 現在のチャンクに行を追加しても最大長を超えない場合
        if len(current_chunk) + len(line) + 1 <= max_length:
            current_chunk += (line + '\n')
            # ブレークポイントを見つけた場合、その位置を記録
            if line.startswith('#'):
                last_breakpoint = len(current_chunk)
        else:
            # 最大長を超える場合の処理
            if last_breakpoint > 0:
                # 直前のブレークポイントまでをチャンクとして追加
                chunks.append(current_chunk[:last_breakpoint])
                current_chunk = current_chunk[last_breakpoint:]
                last_breakpoint = 0
            else:
                # ブレークポイントがない場合は現在のチャンクをそのまま追加
                chunks.append(current_chunk)
                current_chunk = ''
            # 新しい行を現在のチャンクに追加
            current_chunk += (line + '\n')

    # 残りのテキストがあれば最後のチャンクとして追加
    if current_chunk:
        chunks.append(current_chunk)

    return chunks



async def write_to_file(filename: str, text: str) -> None:
    """Asynchronously write text to a file in UTF-8 encoding.

    Args:
        filename (str): The filename to write to.
        text (str): The text to write.
    """
    # Convert text to UTF-8, replacing any problematic characters
    text_utf8 = text.encode('utf-8', errors='replace').decode('utf-8')

    async with aiofiles.open(filename, "w", encoding='utf-8') as file:
        await file.write(text_utf8)


async def write_text_to_md(text: str, path: str) -> str:
    """Writes text to a Markdown file and returns the file path.

    Args:
        text (str): Text to write to the Markdown file.

    Returns:
        str: The file path of the generated Markdown file.
    """
    task = uuid.uuid4().hex
    file_path = f"{path}/{task}.md"
    await write_to_file(file_path, text)
    print(f"Report written to {file_path}")
    return file_path


async def write_md_to_pdf(text: str, path: str) -> str:
    """Converts Markdown text to a PDF file and returns the file path.

    Args:
        text (str): Markdown text to convert.

    Returns:
        str: The encoded file path of the generated PDF.
    """
    task = uuid.uuid4().hex
    file_path = f"{path}/{task}.pdf"
    pdf_styles_path = os.path.join(os.getcwd(), 'app', 'chainlit', 'auto_documentor', 'utils', 'pdf_styles.css')

    try:
        md2pdf(file_path,
               md_content=text,
               # md_file_path=f"{file_path}.md",
               css_file_path=pdf_styles_path,
               base_url=None)
        print(f"Report written to {file_path}")
    except Exception as e:
        print(f"Error in converting Markdown to PDF: {e}")
        return ""

    encoded_file_path = urllib.parse.quote(file_path)
    return encoded_file_path


async def write_md_to_word(text: str, path: str) -> str:
    """Converts Markdown text to a DOCX file and returns the file path.

    Args:
        text (str): Markdown text to convert.

    Returns:
        str: The encoded file path of the generated DOCX.
    """
    task = uuid.uuid4().hex
    file_path = f"{path}/{task}.docx"

    try:
        # Convert report markdown to HTML
        html = mistune.html(text)
        # Create a document object
        doc = Document()
        # Convert the html generated from the report to document format
        HtmlToDocx().add_html_to_document(html, doc)

        # Saving the docx document to file_path
        doc.save(file_path)

        print(f"Report written to {file_path}")

        encoded_file_path = urllib.parse.quote(f"{file_path}.docx")
        return encoded_file_path

    except Exception as e:
        print(f"Error in converting Markdown to DOCX: {e}")
        return ""


async def write_md_to_ppt(text: str, path: str) -> str:
    """Converts Markdown text to a PPTX file and returns the file path.

    Args:
        text (str): Markdown text to convert.

    Returns:
        str: The encoded file path of the generated PPTX.
    """
    task = uuid.uuid4().hex
    file_path = f"{path}/{task}.pptx"


    def generate_slides(marp_content, output_path):
        try:
            # 最初の '---' を削除
            if marp_content.startswith('---\n'):
                marp_content = marp_content[4:]

            # コンテンツをセクションごとに分割
            sections = re.split(r'\n---\n', marp_content.strip())

            # メタデータを解析
            metadata_section = sections[0]
            metadata_lines = metadata_section.strip().split('\n')
            metadata = {}

            for line in metadata_lines:
                if ':' in line:
                    key, value = line.split(':', 1)
                    metadata[key.strip()] = value.strip()

            # プレゼンテーションを作成
            prs = Presentation()

            # スライドサイズを16:9に設定
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # メタデータを設定
            prs.core_properties.title = metadata.get('title', '')
            prs.core_properties.author = metadata.get('author', '')
            date_str = metadata.get('date', '')
            try:
                prs.core_properties.created = datetime.strptime(date_str, '%d/%m/%Y')
            except ValueError:
                pass  # 日付の形式が異なる場合は無視

            # スライドのレイアウトを定義
            blank_slide_layout = prs.slide_layouts[6]  # 'Blank' レイアウト
            title_slide_layout = prs.slide_layouts[0]  # 'Title' レイアウト

            # スライド内容を解析する関数
            def parse_slide_content(slide_content):
                lines = slide_content.strip().split('\n')
                main_heading = None
                blocks = []
                current_block = None
                i = 0
                while i < len(lines):
                    line = lines[i].rstrip()
                    if not line:
                        i += 1
                        continue
                    elif re.match(r'^\s*#\s+', line):
                        # メイン見出し（タイトルスライド）
                        text = line.lstrip('#').strip()
                        main_heading = text
                        i += 1
                        # 日付が次の行にある場合
                        if i < len(lines) and re.match(r'^\s*###\s+', lines[i]):
                            date_text = lines[i].lstrip('#').strip()
                            main_heading += '\n' + date_text
                            i += 1
                    elif re.match(r'^\s*##\s+', line):
                        # スライドのメイン見出し
                        text = line.lstrip('#').strip()
                        main_heading = text
                        i += 1
                    elif re.match(r'^\s*###\s+', line):
                        # サブ見出し（新しいブロックの開始）
                        text = line.lstrip('#').strip()
                        current_block = {'subheading': text, 'bullets': []}
                        blocks.append(current_block)
                        i += 1
                    elif re.match(r'^\s*\*\s+', line):
                        # 箇条書き
                        text = line.lstrip('*').strip()
                        if current_block is not None:
                            current_block['bullets'].append(text)
                        else:
                            # サブ見出しなしで箇条書きが始まった場合
                            current_block = {'subheading': '', 'bullets': [text]}
                            blocks.append(current_block)
                        i += 1
                    else:
                        # その他のテキストは無視
                        i += 1
                return main_heading, blocks

            # 各スライドを処理
            slides_sections = sections[1:]

            for idx, slide_content in enumerate(slides_sections):
                # スライド内容を解析
                main_heading, blocks = parse_slide_content(slide_content)
                
                # メインタイトルスライドの処理
                if idx == 0 and main_heading and not blocks:
                    # タイトルスライドを作成
                    slide = prs.slides.add_slide(title_slide_layout)
                    title_placeholder = slide.shapes.title
                    subtitle_placeholder = slide.placeholders[1]
                    
                    # タイトルを設定
                    if '\n' in main_heading:
                        # 日付が含まれている場合
                        title_text, date_text = main_heading.split('\n', 1)
                        title_placeholder.text = title_text
                        subtitle_placeholder.text = date_text
                    else:
                        title_placeholder.text = main_heading
                        subtitle_placeholder.text = ''
                    
                    # プレースホルダーの位置とサイズを調整
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = prs.slide_width - Inches(1)
                    height = Inches(2)
                    
                    # タイトルの位置とサイズを設定
                    title_placeholder.left = int(left)
                    title_placeholder.top = int(top)
                    title_placeholder.width = int(width)
                    title_placeholder.height = int(height)
                    
                    # サブタイトルの位置とサイズを設定
                    subtitle_placeholder.left = int(left)
                    subtitle_placeholder.top = int(top + title_placeholder.height)
                    subtitle_placeholder.width = int(width)
                    subtitle_placeholder.height = int(Inches(1))
                    
                    # フォントと色を設定
                    for shape in [title_placeholder, subtitle_placeholder]:
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.LEFT  # 左寄せ
                            for run in paragraph.runs:
                                run.font.name = 'Yu Gothic UI'
                                run.font.color.rgb = RGBColor(76, 73, 72)
                    title_placeholder.text_frame.paragraphs[0].font.size = Pt(40)
                    subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
                else:
                    # 通常のスライドを作成
                    slide_layout = prs.slide_layouts[1]  # 'Title and Content' レイアウト
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # タイトルを設定
                    if main_heading:
                        title_placeholder = slide.shapes.title
                        title_placeholder.text = main_heading
                        title_paragraph = title_placeholder.text_frame.paragraphs[0]
                        title_paragraph.alignment = PP_ALIGN.LEFT  # 左寄せ
                        title_paragraph.font.size = Pt(32)
                        title_paragraph.font.bold = True
                        title_paragraph.font.name = 'Yu Gothic UI'
                        title_paragraph.font.color.rgb = RGBColor(76, 73, 72)
                        
                        # タイトルプレースホルダーの位置とサイズを調整
                        left = Inches(0.5)
                        top = Inches(0.5)
                        width = prs.slide_width - Inches(1)
                        height = Inches(1)
                        title_placeholder.left = int(left)
                        title_placeholder.top = int(top)
                        title_placeholder.width = int(width)
                        title_placeholder.height = int(height)
                    
                    # コンテンツプレースホルダーを取得
                    content_placeholder = slide.placeholders[1]
                    
                    # コンテンツプレースホルダーの位置とサイズを調整
                    top_content = top + height + Inches(0.2)
                    height_content = prs.slide_height - top_content - Inches(0.5)
                    content_placeholder.left = int(left)
                    content_placeholder.top = int(top_content)
                    content_placeholder.width = int(width)
                    content_placeholder.height = int(height_content)
                    
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for block in blocks:
                        # サブ見出しを追加
                        p = text_frame.add_paragraph()
                        p.text = block['subheading']
                        p.font.size = Pt(24)
                        p.font.bold = True
                        p.font.name = 'Yu Gothic UI'
                        p.font.color.rgb = RGBColor(76, 73, 72)
                        p.alignment = PP_ALIGN.LEFT  # 左寄せ
                        p.level = 0  # 第一階層
                        
                        # 箇条書きを追加
                        for bullet_text in block['bullets']:
                            bullet_p = text_frame.add_paragraph()
                            bullet_p.text = bullet_text
                            bullet_p.font.size = Pt(18)
                            bullet_p.font.name = 'Yu Gothic UI'
                            bullet_p.font.color.rgb = RGBColor(76, 73, 72)
                            bullet_p.alignment = PP_ALIGN.LEFT  # 左寄せ
                            bullet_p.level = 1  # 箇条書きはサブ見出しよりインデントを深く
                            bullet_p.font.bold = False
                                
                        # 空行を追加
                        p = text_frame.add_paragraph()
                        p.text = ''
                        p.level = 0
                    
                    # 最後の空行を削除
                    if text_frame.paragraphs and text_frame.paragraphs[-1].text == '':
                        text_frame._element.remove(text_frame.paragraphs[-1]._p)

            # プレゼンテーションを保存
            prs.save(output_path)

        except Exception as e:
            print(f"Failed to generate slides: {e}")

    try:

        generate_slides(text, file_path)

        print(f"Report written to {file_path}")

        encoded_file_path = urllib.parse.quote(file_path)   # 拡張子が重複していたので修正
        return encoded_file_path

    except Exception as e:
        print(f"Error in converting Markdown to PPTX: {e}")
        return ""
